"""Tests for file processors."""

import os
import sqlite3
import struct
import zipfile
from unittest.mock import patch

import pytest

from file_processors import (
    CsvProcessor,
    DocxProcessor,
    EmlProcessor,
    HtmlProcessor,
    IcalProcessor,
    JsonProcessor,
    MarkdownProcessor,
    MboxProcessor,
    MsgProcessor,
    OdsProcessor,
    OdtProcessor,
    PdfProcessor,
    PptProcessor,
    PptxProcessor,
    PropertiesProcessor,
    RtfProcessor,
    SqliteProcessor,
    TextProcessor,
    VcfProcessor,
    XlsProcessor,
    XlsxProcessor,
    XmlProcessor,
    YamlProcessor,
    ZipProcessor,
)


class TestPdfProcessor:
    """Tests for PDF processor."""

    def test_can_process_pdf(self):
        """Test that PDF processor recognizes .pdf extension."""
        processor = PdfProcessor()
        assert processor.can_process(".pdf")
        assert processor.can_process(".PDF")
        assert not processor.can_process(".docx")
        assert not processor.can_process(".txt")

    def test_can_process_case_insensitive(self):
        """Test that extension matching is case insensitive."""
        processor = PdfProcessor()
        assert processor.can_process(".PDF")
        assert processor.can_process(".Pdf")

    def test_finalize_page_keeps_short_values_via_accumulation(self):
        """Short standalone values survive because the whole page is yielded together.

        Previously each text container was filtered independently, so a short line
        like a postal code ('50667') in its own container was dropped.  Accumulating
        the page first keeps it whenever the page as a whole has content.
        """
        page = "Kunde Max Mustermann\n50667 Koeln"
        result = PdfProcessor._finalize_page(page, ocr_callable=None)
        assert result == page
        assert "50667" in result

    def test_finalize_page_empty_without_ocr(self):
        """An (almost) empty page yields nothing when OCR is unavailable."""
        assert PdfProcessor._finalize_page("  \n ", ocr_callable=None) == ""
        assert PdfProcessor._finalize_page("x", ocr_callable=None) == ""

    def test_finalize_page_uses_ocr_fallback(self):
        """When a page has no embedded text, the OCR callable result is used."""
        result = PdfProcessor._finalize_page(
            "", ocr_callable=lambda: "Gescannter Text DE89 3704 0044 0532 0130 00"
        )
        assert "DE89 3704 0044 0532 0130 00" in result

    def test_finalize_page_prefers_embedded_text_over_ocr(self):
        """OCR is not invoked when the page already has enough embedded text."""
        calls = []

        def ocr():
            calls.append(1)
            return "should-not-be-used"

        page = "Eingebetteter Text auf dieser Seite"
        assert PdfProcessor._finalize_page(page, ocr_callable=ocr) == page
        assert calls == []  # OCR must not run when embedded text is sufficient

    def test_ocr_available_returns_bool(self):
        """The OCR availability probe never raises and returns a bool."""
        from file_processors.pdf_processor import _ocr_available

        assert isinstance(_ocr_available(), bool)

    def test_ocr_page_passes_language_dpi_and_grayscale(self):
        """_ocr_page forwards configured DPI, greyscale and language to the OCR stack."""
        import sys
        import types

        from core import constants
        from file_processors import pdf_processor

        captured = {}

        def fake_convert_from_path(path, first_page, last_page, dpi, grayscale):
            captured["dpi"] = dpi
            captured["grayscale"] = grayscale
            return ["img1"]

        def fake_image_to_string(img, lang):
            captured["lang"] = lang
            return "Gescannter Text"

        fake_pdf2image = types.ModuleType("pdf2image")
        fake_pdf2image.convert_from_path = fake_convert_from_path
        fake_pytesseract = types.ModuleType("pytesseract")
        fake_pytesseract.image_to_string = fake_image_to_string

        with patch.dict(
            sys.modules,
            {"pdf2image": fake_pdf2image, "pytesseract": fake_pytesseract},
        ):
            with patch.dict(os.environ, {}, clear=False):
                os.environ.pop("PBD_OCR_LANG", None)
                os.environ.pop("PBD_OCR_DPI", None)
                result = pdf_processor._ocr_page("scan.pdf", 1)

        assert result == "Gescannter Text"
        assert captured["dpi"] == constants.OCR_DPI
        assert captured["grayscale"] == constants.OCR_GRAYSCALE
        assert captured["lang"] == constants.OCR_LANGUAGES

    def test_ocr_page_env_overrides(self):
        """PBD_OCR_LANG and PBD_OCR_DPI override the defaults."""
        import sys
        import types

        from file_processors import pdf_processor

        captured = {}

        def fake_convert_from_path(path, first_page, last_page, dpi, grayscale):
            captured["dpi"] = dpi
            return ["img1"]

        def fake_image_to_string(img, lang):
            captured["lang"] = lang
            return "text"

        fake_pdf2image = types.ModuleType("pdf2image")
        fake_pdf2image.convert_from_path = fake_convert_from_path
        fake_pytesseract = types.ModuleType("pytesseract")
        fake_pytesseract.image_to_string = fake_image_to_string

        with patch.dict(
            sys.modules,
            {"pdf2image": fake_pdf2image, "pytesseract": fake_pytesseract},
        ):
            with patch.dict(
                os.environ, {"PBD_OCR_LANG": "fra", "PBD_OCR_DPI": "150"}, clear=False
            ):
                pdf_processor._ocr_page("scan.pdf", 1)

        assert captured["lang"] == "fra"
        assert captured["dpi"] == 150


class TestDocxProcessor:
    """Tests for DOCX processor."""

    def test_can_process_docx(self):
        """Test that DOCX processor recognizes .docx extension."""
        processor = DocxProcessor()
        assert processor.can_process(".docx")
        assert processor.can_process(".DOCX")
        assert not processor.can_process(".pdf")
        assert not processor.can_process(".txt")

    def test_extract_paragraphs_tables_headers(self, temp_dir):
        """Paragraphs, table cells and header/footer text are all extracted."""
        docx = pytest.importorskip("docx")
        path = os.path.join(temp_dir, "sample.docx")

        document = docx.Document()
        document.add_paragraph("Max Mustermann")
        document.add_paragraph("max.mustermann@example.com")
        table = document.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "IBAN"
        table.cell(0, 1).text = "Wert"
        table.cell(1, 0).text = "DE89 3704 0044 0532 0130 00"
        table.cell(1, 1).text = "Hauptkonto"
        section = document.sections[0]
        section.header.paragraphs[0].text = "Vertraulich Kopfzeile"
        section.footer.paragraphs[0].text = "Seite Fusszeile"
        document.save(path)

        text = DocxProcessor().extract_text(path)

        # Paragraphs are newline-separated so adjacent entities do not fuse.
        assert "Max Mustermann\nmax.mustermann@example.com" in text
        # Table cells (incl. the IBAN that only lives in a table) are present.
        assert "DE89 3704 0044 0532 0130 00" in text
        assert "IBAN" in text
        # Header and footer text is captured.
        assert "Vertraulich Kopfzeile" in text
        assert "Seite Fusszeile" in text


class TestHtmlProcessor:
    """Tests for HTML processor."""

    def test_can_process_html(self):
        """Test that HTML processor recognizes .html extension."""
        processor = HtmlProcessor()
        assert processor.can_process(".html")
        assert processor.can_process(".HTML")
        assert processor.can_process(".htm")
        assert not processor.can_process(".txt")

    def test_extract_text_from_html(self, sample_html_file):
        """Test text extraction from HTML file."""
        processor = HtmlProcessor()
        text = processor.extract_text(sample_html_file)
        assert "user@example.com" in text
        assert "IBAN" in text
        # HTML tags should be removed
        assert "<html>" not in text
        assert "<p>" not in text
        assert "<" not in text and ">" not in text


class TestTextProcessor:
    """Tests for text processor."""

    def test_can_process_txt(self, temp_dir):
        """Test that text processor recognizes .txt extension."""
        processor = TextProcessor()
        assert processor.can_process(".txt", os.path.join(temp_dir, "test.txt"))
        assert processor.can_process(".TXT", os.path.join(temp_dir, "test.TXT"))
        assert not processor.can_process(".pdf", os.path.join(temp_dir, "test.pdf"))

    def test_extract_text_from_file(self, sample_text_file):
        """Test text extraction from text file."""
        processor = TextProcessor()
        text = processor.extract_text(sample_text_file)
        assert "test@example.com" in text
        assert "IBAN" in text
        assert "This is a test file" in text

    def test_file_not_found(self, temp_dir):
        """Test that FileNotFoundError is raised for non-existent file."""
        processor = TextProcessor()
        non_existent = os.path.join(temp_dir, "nonexistent.txt")
        with pytest.raises(FileNotFoundError):
            processor.extract_text(non_existent)


class TestCsvProcessor:
    """Tests for CSV processor."""

    def test_can_process_csv(self):
        """Test that CSV processor recognizes .csv extension."""
        processor = CsvProcessor()
        assert processor.can_process(".csv")
        assert processor.can_process(".CSV")
        assert not processor.can_process(".txt")
        assert not processor.can_process(".json")

    def test_extract_text_from_csv(self, temp_dir):
        """Test text extraction from CSV file."""
        file_path = os.path.join(temp_dir, "test.csv")
        with open(file_path, "w", encoding="utf-8") as f:
            f.write("Name,Email,Phone\n")
            f.write("John Doe,john@example.com,123-456-7890\n")
            f.write("Jane Smith,jane@example.com,098-765-4321\n")

        processor = CsvProcessor()
        text = processor.extract_text(file_path)
        assert "John Doe" in text
        assert "john@example.com" in text
        assert "Jane Smith" in text
        assert "jane@example.com" in text
        assert "123-456-7890" in text

    def test_extract_text_from_csv_semicolon(self, temp_dir):
        """Test text extraction from CSV file with semicolon delimiter."""
        file_path = os.path.join(temp_dir, "test.csv")
        with open(file_path, "w", encoding="utf-8") as f:
            f.write("Name;Email;Phone\n")
            f.write("John Doe;john@example.com;123-456-7890\n")

        processor = CsvProcessor()
        text = processor.extract_text(file_path)
        assert "John Doe" in text
        assert "john@example.com" in text

    def test_extract_text_from_csv_tab(self, temp_dir):
        """Test text extraction from CSV file with tab delimiter."""
        file_path = os.path.join(temp_dir, "test.csv")
        with open(file_path, "w", encoding="utf-8") as f:
            f.write("Name\tEmail\tPhone\n")
            f.write("John Doe\tjohn@example.com\t123-456-7890\n")

        processor = CsvProcessor()
        text = processor.extract_text(file_path)
        assert "John Doe" in text
        assert "john@example.com" in text

    def test_extract_preserves_column_context(self, temp_dir):
        """Each value is paired with its column header and rows stay separated."""
        file_path = os.path.join(temp_dir, "context.csv")
        with open(file_path, "w", encoding="utf-8") as f:
            f.write("Name,IBAN\n")
            f.write("Max Mustermann,DE89 3704 0044 0532 0130 00\n")
            f.write("Erika Beispiel,DE02 1203 0000 0000 2020 51\n")

        text = CsvProcessor().extract_text(file_path)
        assert "IBAN: DE89 3704 0044 0532 0130 00" in text
        assert "Name: Max Mustermann" in text
        # Two data records -> two separate lines (plus the header line).
        assert len(text.splitlines()) == 3

    def test_file_not_found(self, temp_dir):
        """Test that FileNotFoundError is raised for non-existent file."""
        processor = CsvProcessor()
        non_existent = os.path.join(temp_dir, "nonexistent.csv")
        with pytest.raises(FileNotFoundError):
            processor.extract_text(non_existent)


class TestJsonProcessor:
    """Tests for JSON processor."""

    def test_can_process_json(self):
        """Test that JSON processor recognizes .json extension."""
        processor = JsonProcessor()
        assert processor.can_process(".json")
        assert processor.can_process(".JSON")
        assert not processor.can_process(".txt")
        assert not processor.can_process(".csv")

    def test_extract_text_from_json(self, temp_dir):
        """Test text extraction from JSON file."""
        file_path = os.path.join(temp_dir, "test.json")
        with open(file_path, "w", encoding="utf-8") as f:
            f.write(
                '{"name": "John Doe", "email": "john@example.com", "phone": "123-456-7890"}'
            )

        processor = JsonProcessor()
        text = processor.extract_text(file_path)
        assert "name" in text
        assert "John Doe" in text
        assert "email" in text
        assert "john@example.com" in text
        assert "phone" in text
        assert "123-456-7890" in text

    def test_extract_text_from_nested_json(self, temp_dir):
        """Test text extraction from nested JSON file."""
        file_path = os.path.join(temp_dir, "test.json")
        with open(file_path, "w", encoding="utf-8") as f:
            f.write(
                """{
                "users": [
                    {"name": "John Doe", "email": "john@example.com"},
                    {"name": "Jane Smith", "email": "jane@example.com"}
                ],
                "metadata": {
                    "created": "2024-01-01",
                    "author": "Admin"
                }
            }"""
            )

        processor = JsonProcessor()
        text = processor.extract_text(file_path)
        assert "users" in text
        assert "John Doe" in text
        assert "john@example.com" in text
        assert "Jane Smith" in text
        assert "jane@example.com" in text
        assert "metadata" in text
        assert "created" in text
        assert "author" in text
        assert "Admin" in text

    def test_extract_text_from_json_array(self, temp_dir):
        """Test text extraction from JSON array."""
        file_path = os.path.join(temp_dir, "test.json")
        with open(file_path, "w", encoding="utf-8") as f:
            f.write('["John Doe", "jane@example.com", "123 Main St"]')

        processor = JsonProcessor()
        text = processor.extract_text(file_path)
        assert "John Doe" in text
        assert "jane@example.com" in text
        assert "123 Main St" in text

    def test_extract_text_from_invalid_json(self, temp_dir):
        """Test text extraction from invalid JSON file (should still extract strings)."""
        file_path = os.path.join(temp_dir, "test.json")
        with open(file_path, "w", encoding="utf-8") as f:
            f.write('{"name": "John Doe", "email": "john@example.com" invalid json')

        processor = JsonProcessor()
        text = processor.extract_text(file_path)
        # Should still extract string values using regex fallback
        assert "John Doe" in text or "john@example.com" in text

    def test_file_not_found(self, temp_dir):
        """Test that FileNotFoundError is raised for non-existent file."""
        processor = JsonProcessor()
        non_existent = os.path.join(temp_dir, "nonexistent.json")
        with pytest.raises(FileNotFoundError):
            processor.extract_text(non_existent)


class TestRtfProcessor:
    """Tests for RTF processor."""

    def test_can_process_rtf(self):
        """Test that RTF processor recognizes .rtf extension."""
        processor = RtfProcessor()
        assert processor.can_process(".rtf")
        assert processor.can_process(".RTF")
        assert not processor.can_process(".txt")
        assert not processor.can_process(".docx")

    def test_extract_text_from_rtf(self, temp_dir):
        """Test text extraction from RTF file."""
        file_path = os.path.join(temp_dir, "test.rtf")
        # Create a simple RTF file
        rtf_content = r"{\rtf1\ansi\deff0 {\fonttbl {\f0 Times New Roman;}}\f0\fs24 This is a test RTF document with email test@example.com and IBAN DE89 3704 0044 0532 0130 00.}"
        with open(file_path, "w", encoding="utf-8") as f:
            f.write(rtf_content)

        processor = RtfProcessor()
        text = processor.extract_text(file_path)
        assert "test@example.com" in text
        assert "IBAN" in text
        assert "test rtf document" in text.lower()

    def test_file_not_found(self, temp_dir):
        """Test that FileNotFoundError is raised for non-existent file."""
        processor = RtfProcessor()
        non_existent = os.path.join(temp_dir, "nonexistent.rtf")
        with pytest.raises(FileNotFoundError):
            processor.extract_text(non_existent)


class TestOdtProcessor:
    """Tests for ODT processor."""

    def test_can_process_odt(self):
        """Test that ODT processor recognizes .odt extension."""
        processor = OdtProcessor()
        assert processor.can_process(".odt")
        assert processor.can_process(".ODT")
        assert not processor.can_process(".txt")
        assert not processor.can_process(".docx")

    def test_file_not_found(self, temp_dir):
        """Test that FileNotFoundError is raised for non-existent file."""
        pytest.importorskip("odf.opendocument")
        processor = OdtProcessor()
        non_existent = os.path.join(temp_dir, "nonexistent.odt")
        with pytest.raises(FileNotFoundError):
            processor.extract_text(non_existent)

    def test_extract_text_from_odt(self, temp_dir):
        """Test text extraction from ODT file (requires odfpy)."""
        pytest.importorskip("odf.opendocument")
        from odf.opendocument import OpenDocumentText
        from odf.text import P

        odt_path = os.path.join(temp_dir, "test.odt")
        doc = OpenDocumentText()
        doc.text.addElement(P(text="Contact: John Doe at john@example.com"))
        doc.save(odt_path)

        processor = OdtProcessor()
        text = processor.extract_text(odt_path)
        assert "John Doe" in text
        assert "john@example.com" in text


class TestEmlProcessor:
    """Tests for EML processor."""

    def test_can_process_eml(self):
        """Test that EML processor recognizes .eml extension."""
        processor = EmlProcessor()
        assert processor.can_process(".eml")
        assert processor.can_process(".EML")
        assert not processor.can_process(".txt")
        assert not processor.can_process(".msg")

    def test_extract_text_from_attachment(self, temp_dir):
        """Text inside a CSV attachment is extracted via the processor registry."""
        from email.message import EmailMessage

        msg = EmailMessage()
        msg["From"] = "alice@example.com"
        msg["To"] = "bob@example.com"
        msg["Subject"] = "Stammdaten"
        msg.set_content("Anbei die Kundendaten.")
        msg.add_attachment(
            b"Name,IBAN\nMax Mustermann,DE89 3704 0044 0532 0130 00\n",
            maintype="text",
            subtype="csv",
            filename="kunden.csv",
        )

        file_path = os.path.join(temp_dir, "with_attachment.eml")
        with open(file_path, "wb") as f:
            f.write(msg.as_bytes())

        text = EmlProcessor().extract_text(file_path)
        # The IBAN only exists inside the attachment.
        assert "DE89 3704 0044 0532 0130 00" in text
        assert "[Attachment: kunden.csv]" in text
        # Attachment is run through the CSV processor -> column context preserved.
        assert "IBAN: DE89 3704 0044 0532 0130 00" in text

    def test_extract_text_from_eml(self, temp_dir):
        """Test text extraction from EML file."""
        file_path = os.path.join(temp_dir, "test.eml")
        # Create a simple EML file
        eml_content = """From: sender@example.com
To: recipient@example.com
Subject: Test Email
Content-Type: text/plain; charset=utf-8

This is a test email with contact info: contact@example.com
IBAN: DE89 3704 0044 0532 0130 00
"""
        with open(file_path, "w", encoding="utf-8") as f:
            f.write(eml_content)

        processor = EmlProcessor()
        text = processor.extract_text(file_path)
        assert "sender@example.com" in text
        assert "recipient@example.com" in text
        assert "Test Email" in text
        assert "contact@example.com" in text
        assert "IBAN" in text

    def test_extract_text_from_multipart_eml(self, temp_dir):
        """Test text extraction from multipart EML file."""
        file_path = os.path.join(temp_dir, "test.eml")
        # Create a multipart EML file
        eml_content = """From: sender@example.com
To: recipient@example.com
Subject: Multipart Test
MIME-Version: 1.0
Content-Type: multipart/alternative; boundary="boundary123"

--boundary123
Content-Type: text/plain; charset=utf-8

This is the plain text part with email test@example.com

--boundary123
Content-Type: text/html; charset=utf-8

<html><body>This is the <b>HTML</b> part with email test@example.com</body></html>

--boundary123--
"""
        with open(file_path, "w", encoding="utf-8") as f:
            f.write(eml_content)

        processor = EmlProcessor()
        text = processor.extract_text(file_path)
        assert "sender@example.com" in text
        assert "recipient@example.com" in text
        assert "test@example.com" in text

    def test_html_part_with_bogus_charset_falls_back_instead_of_dropping_content(
        self, temp_dir
    ):
        """An unrecognised charset on the HTML part must not silently lose the body.

        Previously a ``LookupError``/``UnicodeDecodeError`` while decoding the
        text/html part was swallowed with a bare ``pass``, dropping the entire
        HTML body. It must now fall back to utf-8, like the text/plain branch.
        """
        file_path = os.path.join(temp_dir, "bad_charset.eml")
        eml_content = (
            "From: sender@example.com\n"
            "To: recipient@example.com\n"
            "Subject: Bad Charset\n"
            "MIME-Version: 1.0\n"
            'Content-Type: multipart/alternative; boundary="boundary123"\n'
            "\n"
            "--boundary123\n"
            'Content-Type: text/html; charset="not-a-real-charset"\n'
            "\n"
            "<html><body>Contact: contact@example.com</body></html>\n"
            "\n"
            "--boundary123--\n"
        )
        with open(file_path, "w", encoding="utf-8") as f:
            f.write(eml_content)

        text = EmlProcessor().extract_text(file_path)
        assert "contact@example.com" in text

    def test_file_not_found(self, temp_dir):
        """Test that FileNotFoundError is raised for non-existent file."""
        processor = EmlProcessor()
        non_existent = os.path.join(temp_dir, "nonexistent.eml")
        with pytest.raises(FileNotFoundError):
            processor.extract_text(non_existent)


class TestMsgProcessor:
    """Tests for MSG processor."""

    def test_can_process_msg(self):
        """Test that MSG processor recognizes .msg extension."""
        processor = MsgProcessor()
        assert processor.can_process(".msg")
        assert processor.can_process(".MSG")
        assert not processor.can_process(".txt")
        assert not processor.can_process(".eml")

    def test_import_error_when_extract_msg_not_installed(self, temp_dir, mocker):
        """Test that ImportError is raised when extract-msg is not installed."""
        # Mock the import to raise ImportError
        mocker.patch(
            "file_processors.msg_processor.extract_msg",
            side_effect=ImportError("No module named 'extract_msg'"),
        )

        processor = MsgProcessor()
        file_path = os.path.join(temp_dir, "test.msg")
        # Create a dummy file (won't be read due to import error)
        with open(file_path, "w") as f:
            f.write("dummy")

        with pytest.raises(ImportError) as exc_info:
            processor.extract_text(file_path)
        assert "extract-msg is required" in str(exc_info.value)

    def test_file_not_found(self, temp_dir):
        """Test that FileNotFoundError is raised for non-existent file."""
        processor = MsgProcessor()
        non_existent = os.path.join(temp_dir, "nonexistent.msg")
        # Note: This will raise ImportError if extract-msg is not installed,
        # or FileNotFoundError if it is installed. We test both cases.
        try:
            with pytest.raises((FileNotFoundError, ImportError)):
                processor.extract_text(non_existent)
        except ImportError:
            # If extract-msg is not installed, that's expected
            pass

    # Note: Testing actual MSG extraction would require creating a valid MSG file
    # which is complex and requires Outlook or specialized tools. The can_process test
    # and error handling tests verify the basic functionality. Full integration tests
    # would require sample MSG files from actual Outlook exports.


class TestOdsProcessor:
    """Tests for ODS processor."""

    def test_can_process_ods(self):
        """Test that ODS processor recognizes .ods extension."""
        processor = OdsProcessor()
        assert processor.can_process(".ods")
        assert processor.can_process(".ODS")
        assert not processor.can_process(".txt")
        assert not processor.can_process(".xlsx")

    def test_file_not_found(self, temp_dir):
        """Test that FileNotFoundError is raised for non-existent file."""
        pytest.importorskip("odf.opendocument")
        processor = OdsProcessor()
        non_existent = os.path.join(temp_dir, "nonexistent.ods")
        with pytest.raises(FileNotFoundError):
            processor.extract_text(non_existent)

    def test_extract_text_from_ods(self, temp_dir):
        """Test text extraction from ODS file (requires odfpy)."""
        pytest.importorskip("odf.opendocument")
        from odf.opendocument import OpenDocumentSpreadsheet
        from odf.table import Table, TableCell, TableRow
        from odf.text import P

        ods_path = os.path.join(temp_dir, "test.ods")
        doc = OpenDocumentSpreadsheet()
        table = Table(name="Sheet1")
        row = TableRow()
        for cell_text in ["John Doe", "john@example.com"]:
            cell = TableCell()
            cell.addElement(P(text=cell_text))
            row.addElement(cell)
        table.addElement(row)
        doc.spreadsheet.addElement(table)
        doc.save(ods_path)

        processor = OdsProcessor()
        text = processor.extract_text(ods_path)
        assert "John Doe" in text
        assert "john@example.com" in text


class TestXlsxProcessor:
    """Tests for XLSX processor."""

    def test_can_process_xlsx(self):
        """Test that XLSX processor recognizes .xlsx extension."""
        processor = XlsxProcessor()
        assert processor.can_process(".xlsx")
        assert processor.can_process(".XLSX")
        assert not processor.can_process(".xls")
        assert not processor.can_process(".csv")

    def test_extract_text_from_xlsx(self, temp_dir):
        """Test text extraction from XLSX file (requires openpyxl)."""
        pytest.importorskip("openpyxl")
        from openpyxl import Workbook

        xlsx_path = os.path.join(temp_dir, "test.xlsx")
        wb = Workbook()
        ws = wb.active
        ws["A1"] = "John Doe"
        ws["B1"] = "john@example.com"
        wb.save(xlsx_path)

        processor = XlsxProcessor()
        text = processor.extract_text(xlsx_path)
        assert "John Doe" in text
        assert "john@example.com" in text

    def test_extract_preserves_column_context(self, temp_dir):
        """Values are paired with their column header for context-aware detection."""
        pytest.importorskip("openpyxl")
        from openpyxl import Workbook

        xlsx_path = os.path.join(temp_dir, "context.xlsx")
        wb = Workbook()
        ws = wb.active
        ws.append(["Name", "IBAN"])
        ws.append(["Max Mustermann", "DE89 3704 0044 0532 0130 00"])
        ws.append(["Erika Beispiel", "DE02 1203 0000 0000 2020 51"])
        wb.save(xlsx_path)

        text = XlsxProcessor().extract_text(xlsx_path)
        # Each value carries its column header.
        assert "IBAN: DE89 3704 0044 0532 0130 00" in text
        assert "Name: Max Mustermann" in text
        # Records stay on separate lines so entities do not fuse across rows.
        assert "Erika Beispiel" in text
        assert "\n" in text

    def test_file_not_found(self, temp_dir):
        """Test that non-existent file raises an error (XlsxProcessor wraps as Exception)."""
        processor = XlsxProcessor()
        non_existent = os.path.join(temp_dir, "nonexistent.xlsx")
        with pytest.raises(Exception) as exc_info:
            processor.extract_text(non_existent)
        assert (
            "No such file" in str(exc_info.value)
            or "nonexistent" in str(exc_info.value).lower()
        )


class TestXlsProcessor:
    """Tests for XLS (Excel 97-2003) processor."""

    def test_can_process_xls(self):
        """Test that XLS processor recognizes .xls extension only."""
        processor = XlsProcessor()
        assert processor.can_process(".xls")
        assert processor.can_process(".XLS")
        assert not processor.can_process(".xlsx")
        assert not processor.can_process(".csv")

    def test_extract_text_from_xls(self, temp_dir):
        """Test text extraction from a real XLS file (requires xlrd + xlwt)."""
        pytest.importorskip("xlrd")
        xlwt = pytest.importorskip("xlwt")

        xls_path = os.path.join(temp_dir, "test.xls")
        wb = xlwt.Workbook()
        ws = wb.add_sheet("Sheet1")
        ws.write(0, 0, "John Doe")
        ws.write(0, 1, "john@example.com")
        wb.save(xls_path)

        processor = XlsProcessor()
        text = processor.extract_text(xls_path)
        assert "John Doe" in text
        assert "john@example.com" in text

    def test_extract_preserves_column_context(self, temp_dir):
        """Values are paired with their column header, same as XlsxProcessor."""
        pytest.importorskip("xlrd")
        xlwt = pytest.importorskip("xlwt")

        xls_path = os.path.join(temp_dir, "context.xls")
        wb = xlwt.Workbook()
        ws = wb.add_sheet("Sheet1")
        rows = [
            ["Name", "IBAN"],
            ["Max Mustermann", "DE89 3704 0044 0532 0130 00"],
            ["Erika Beispiel", "DE02 1203 0000 0000 2020 51"],
        ]
        for r, row in enumerate(rows):
            for c, value in enumerate(row):
                ws.write(r, c, value)
        wb.save(xls_path)

        text = XlsProcessor().extract_text(xls_path)
        assert "IBAN: DE89 3704 0044 0532 0130 00" in text
        assert "Name: Max Mustermann" in text
        assert "Erika Beispiel" in text

    def test_import_error_when_xlrd_not_installed(self, temp_dir, mocker):
        """Test that ImportError is raised when xlrd is not installed."""
        mocker.patch.dict("sys.modules", {"xlrd": None})

        processor = XlsProcessor()
        file_path = os.path.join(temp_dir, "test.xls")
        with open(file_path, "w") as f:
            f.write("dummy")

        with pytest.raises(ImportError) as exc_info:
            processor.extract_text(file_path)
        assert "xlrd is required" in str(exc_info.value)

    def test_file_not_found(self, temp_dir):
        """Test that non-existent file raises an error (XlsProcessor wraps as Exception)."""
        pytest.importorskip("xlrd")
        processor = XlsProcessor()
        non_existent = os.path.join(temp_dir, "nonexistent.xls")
        with pytest.raises(Exception) as exc_info:
            processor.extract_text(non_existent)
        assert (
            "No such file" in str(exc_info.value)
            or "nonexistent" in str(exc_info.value).lower()
        )


class TestPptxProcessor:
    """Tests for PPTX processor."""

    def test_can_process_pptx(self):
        """Test that PPTX processor recognizes .pptx extension."""
        processor = PptxProcessor()
        assert processor.can_process(".pptx")
        assert processor.can_process(".PPTX")
        assert not processor.can_process(".txt")
        assert not processor.can_process(".docx")

    def test_import_error_when_python_pptx_not_installed(self, temp_dir, mocker):
        """Test that ImportError is raised when python-pptx is not installed."""
        # Mock the import to raise ImportError
        mocker.patch(
            "file_processors.pptx_processor.Presentation",
            side_effect=ImportError("No module named 'pptx'"),
        )

        processor = PptxProcessor()
        file_path = os.path.join(temp_dir, "test.pptx")
        # Create a dummy file (won't be read due to import error)
        with open(file_path, "w") as f:
            f.write("dummy")

        with pytest.raises(ImportError) as exc_info:
            processor.extract_text(file_path)
        assert "python-pptx is required" in str(exc_info.value)

    def test_file_not_found(self, temp_dir):
        """Test that FileNotFoundError is raised for non-existent file."""
        processor = PptxProcessor()
        non_existent = os.path.join(temp_dir, "nonexistent.pptx")
        # Note: This will raise ImportError if python-pptx is not installed,
        # or FileNotFoundError if it is installed. We test both cases.
        try:
            with pytest.raises((FileNotFoundError, ImportError)):
                processor.extract_text(non_existent)
        except ImportError:
            # If python-pptx is not installed, that's expected
            pass

    def test_extract_text_from_pptx(self, temp_dir):
        """Test text extraction from PPTX file (requires python-pptx)."""
        pytest.importorskip("pptx")
        from pptx import Presentation
        from pptx.util import Inches

        pptx_path = os.path.join(temp_dir, "test.pptx")
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(5), Inches(1)
        ).text_frame.text = "John Doe john@example.com"
        prs.save(pptx_path)

        processor = PptxProcessor()
        text = processor.extract_text(pptx_path)
        assert "John Doe" in text
        assert "john@example.com" in text


def _build_minimal_ppt(texts: list[str]) -> bytes:
    """Hand-build a minimal, real OLE2 (CFBF) file with a "PowerPoint Document"
    stream containing the given strings as ``TextCharsAtom`` records.

    There is no maintained pure-Python library that *writes* the legacy PPT
    binary format (that gap is exactly what ``PptProcessor`` closes for
    *reading*), and LibreOffice is not available for fixture generation in CI.
    So this builds the smallest valid compound-file container by hand: one
    FAT sector, one directory sector (Root Entry + the document stream), and
    enough data sectors to hold the record payload -- deliberately avoiding
    the mini-stream/mini-FAT mechanism by padding the stream to at least 4096
    bytes (the standard mini-stream cutoff) so only regular sectors are used.
    This keeps the layout uniform regardless of how much text is packed in,
    at the cost of only handling small payloads (asserted below): fine for a
    test fixture, not a general-purpose OLE writer.
    """
    NOSTREAM = 0xFFFFFFFF
    ENDOFCHAIN = 0xFFFFFFFE
    FREESECT = 0xFFFFFFFF
    FATSECT = 0xFFFFFFFD
    TEXT_CHARS_ATOM = 4008

    def atom(rec_type: int, content: bytes) -> bytes:
        return struct.pack("<HHI", 0x0, rec_type, len(content)) + content

    def container(rec_type: int, children: list[bytes]) -> bytes:
        content = b"".join(children)
        return struct.pack("<HHI", 0xF, rec_type, len(content)) + content

    payload = container(
        1000, [atom(TEXT_CHARS_ATOM, t.encode("utf-16-le")) for t in texts]
    )
    data_size = max(4096, -(-len(payload) // 512) * 512)
    stream = payload + b"\x00" * (data_size - len(payload))
    num_data_sectors = data_size // 512
    dir_sector_index = 1 + num_data_sectors
    total_sectors = 1 + num_data_sectors + 1  # FAT sector + data + directory
    assert total_sectors <= 128, "fixture text too large for this minimal builder"

    def name_field(name: str) -> tuple[bytes, int]:
        raw = name.encode("utf-16-le") + b"\x00\x00"
        return raw + b"\x00" * (64 - len(raw)), len(raw)

    def dir_entry(
        name: str,
        obj_type: int,
        left: int,
        right: int,
        child: int,
        start_sector: int,
        stream_size: int,
    ) -> bytes:
        name_bytes, name_len = name_field(name)
        return (
            name_bytes
            + struct.pack("<H", name_len)
            + struct.pack("<B", obj_type)
            + struct.pack("<B", 1)  # color flag (black)
            + struct.pack("<I", left)
            + struct.pack("<I", right)
            + struct.pack("<I", child)
            + b"\x00" * 16  # CLSID
            + struct.pack("<I", 0)  # state bits
            + b"\x00" * 16  # creation + modified time
            + struct.pack("<I", start_sector)
            + struct.pack("<Q", stream_size)
        )

    root_entry = dir_entry("Root Entry", 5, NOSTREAM, NOSTREAM, 1, ENDOFCHAIN, 0)
    stream_entry = dir_entry(
        "PowerPoint Document", 2, NOSTREAM, NOSTREAM, NOSTREAM, 1, data_size
    )
    dir_sector = root_entry + stream_entry + b"\x00" * 128 + b"\x00" * 128

    fat_entries = [FATSECT]
    fat_entries += [s + 1 for s in range(1, num_data_sectors)]
    fat_entries += [ENDOFCHAIN]  # last data sector
    fat_entries += [ENDOFCHAIN]  # directory sector (single sector)
    fat_entries += [FREESECT] * (128 - len(fat_entries))
    fat_sector = b"".join(struct.pack("<I", e) for e in fat_entries)

    difat = [0] + [FREESECT] * 108  # only FAT sector 0 is in use

    header = (
        bytes.fromhex("d0cf11e0a1b11ae1")  # OLE2 signature
        + b"\x00" * 16  # CLSID
        + struct.pack("<H", 0x003E)  # minor version
        + struct.pack("<H", 0x0003)  # major version (3 = 512-byte sectors)
        + struct.pack("<H", 0xFFFE)  # byte order mark
        + struct.pack("<H", 9)  # sector shift: 2^9 = 512
        + struct.pack("<H", 6)  # mini sector shift: 2^6 = 64
        + b"\x00" * 6  # reserved
        + struct.pack("<I", 0)  # number of directory sectors (0 for v3)
        + struct.pack("<I", 1)  # number of FAT sectors
        + struct.pack("<I", dir_sector_index)  # first directory sector
        + struct.pack("<I", 0)  # transaction signature
        + struct.pack("<I", 4096)  # mini stream cutoff size
        + struct.pack("<I", ENDOFCHAIN)  # first mini FAT sector (none)
        + struct.pack("<I", 0)  # number of mini FAT sectors
        + struct.pack("<I", ENDOFCHAIN)  # first DIFAT sector (none beyond header)
        + struct.pack("<I", 0)  # number of DIFAT sectors
        + b"".join(struct.pack("<I", e) for e in difat)
    )
    return header + fat_sector + stream + dir_sector


class TestPptProcessor:
    """Tests for PPT (PowerPoint 97-2003) processor."""

    def test_can_process_ppt(self):
        """Test that PPT processor recognizes .ppt extension."""
        processor = PptProcessor()
        assert processor.can_process(".ppt")
        assert processor.can_process(".PPT")
        assert not processor.can_process(".txt")
        assert not processor.can_process(".pptx")

    def test_extract_text_from_ppt(self, temp_dir):
        """Test text extraction from a real legacy PPT file (requires olefile)."""
        pytest.importorskip("olefile")

        ppt_path = os.path.join(temp_dir, "test.ppt")
        with open(ppt_path, "wb") as f:
            f.write(_build_minimal_ppt(["John Doe", "john@example.com"]))

        processor = PptProcessor()
        text = processor.extract_text(ppt_path)
        assert "John Doe" in text
        assert "john@example.com" in text

    def test_extract_text_spanning_multiple_sectors(self, temp_dir):
        """Text long enough to span several 512-byte sectors is still fully read."""
        pytest.importorskip("olefile")

        long_run = "A" * 3000
        ppt_path = os.path.join(temp_dir, "large.ppt")
        with open(ppt_path, "wb") as f:
            f.write(_build_minimal_ppt(["John Doe", long_run]))

        text = PptProcessor().extract_text(ppt_path)
        assert "John Doe" in text
        assert long_run in text

    def test_import_error_when_olefile_not_installed(self, temp_dir, mocker):
        """Test that ImportError is raised when olefile is not installed."""
        mocker.patch("file_processors.pptx_processor.olefile", None)

        processor = PptProcessor()
        file_path = os.path.join(temp_dir, "test.ppt")
        with open(file_path, "w") as f:
            f.write("dummy")

        with pytest.raises(ImportError) as exc_info:
            processor.extract_text(file_path)
        assert "olefile is required" in str(exc_info.value)

    def test_not_an_ole_file(self, temp_dir):
        """Test that a non-OLE file (e.g. a renamed text file) raises clearly."""
        pytest.importorskip("olefile")

        file_path = os.path.join(temp_dir, "not_really.ppt")
        with open(file_path, "w") as f:
            f.write("this is not a compound file")

        with pytest.raises(Exception) as exc_info:
            PptProcessor().extract_text(file_path)
        assert "Not a valid OLE compound file" in str(exc_info.value)

    def test_file_not_found(self, temp_dir):
        """Test that FileNotFoundError is raised for a non-existent file."""
        pytest.importorskip("olefile")

        processor = PptProcessor()
        non_existent = os.path.join(temp_dir, "nonexistent.ppt")
        with pytest.raises(FileNotFoundError):
            processor.extract_text(non_existent)


class TestYamlProcessor:
    """Tests for YAML processor."""

    def test_can_process_yaml(self):
        """Test that YAML processor recognizes .yaml and .yml extensions."""
        processor = YamlProcessor()
        assert processor.can_process(".yaml")
        assert processor.can_process(".YAML")
        assert processor.can_process(".yml")
        assert processor.can_process(".YML")
        assert not processor.can_process(".txt")
        assert not processor.can_process(".json")

    def test_extract_text_from_yaml(self, temp_dir):
        """Test text extraction from YAML file."""
        file_path = os.path.join(temp_dir, "test.yaml")
        with open(file_path, "w", encoding="utf-8") as f:
            f.write(
                """name: John Doe
email: john@example.com
phone: "123-456-7890"
address:
  street: "123 Main St"
  city: "New York"
"""
            )

        processor = YamlProcessor()
        text = processor.extract_text(file_path)
        assert "name" in text
        assert "John Doe" in text
        assert "email" in text
        assert "john@example.com" in text
        assert "phone" in text
        assert "123-456-7890" in text
        assert "address" in text
        assert "street" in text
        assert "123 Main St" in text
        assert "city" in text
        assert "New York" in text

    def test_extract_text_from_nested_yaml(self, temp_dir):
        """Test text extraction from nested YAML file."""
        file_path = os.path.join(temp_dir, "test.yaml")
        with open(file_path, "w", encoding="utf-8") as f:
            f.write(
                """users:
  - name: John Doe
    email: john@example.com
  - name: Jane Smith
    email: jane@example.com
metadata:
  created: "2024-01-01"
  author: Admin
"""
            )

        processor = YamlProcessor()
        text = processor.extract_text(file_path)
        assert "users" in text
        assert "John Doe" in text
        assert "john@example.com" in text
        assert "Jane Smith" in text
        assert "jane@example.com" in text
        assert "metadata" in text
        assert "created" in text
        assert "author" in text
        assert "Admin" in text

    def test_extract_text_from_yaml_array(self, temp_dir):
        """Test text extraction from YAML array."""
        file_path = os.path.join(temp_dir, "test.yaml")
        with open(file_path, "w", encoding="utf-8") as f:
            f.write('- "John Doe"\n- "jane@example.com"\n- "123 Main St"')

        processor = YamlProcessor()
        text = processor.extract_text(file_path)
        assert "John Doe" in text
        assert "jane@example.com" in text
        assert "123 Main St" in text

    def test_import_error_when_pyyaml_not_installed(self, temp_dir, mocker):
        """Test that ImportError is raised when PyYAML is not installed."""
        # Mock the import to raise ImportError
        mocker.patch(
            "file_processors.yaml_processor.yaml",
            side_effect=ImportError("No module named 'yaml'"),
        )

        processor = YamlProcessor()
        file_path = os.path.join(temp_dir, "test.yaml")
        with open(file_path, "w") as f:
            f.write("key: value")

        with pytest.raises(ImportError) as exc_info:
            processor.extract_text(file_path)
        assert "PyYAML is required" in str(exc_info.value)

    def test_file_not_found(self, temp_dir):
        """Test that FileNotFoundError is raised for non-existent file."""
        processor = YamlProcessor()
        non_existent = os.path.join(temp_dir, "nonexistent.yaml")
        # Note: This will raise ImportError if PyYAML is not installed,
        # or FileNotFoundError if it is installed. We test both cases.
        try:
            with pytest.raises((FileNotFoundError, ImportError)):
                processor.extract_text(non_existent)
        except ImportError:
            # If PyYAML is not installed, that's expected
            pass


class TestMarkdownProcessor:
    """Tests for Markdown processor."""

    def test_can_process_md(self):
        """Test that Markdown processor recognizes .md extension."""
        processor = MarkdownProcessor()
        assert processor.can_process(".md")
        assert processor.can_process(".MD")
        assert processor.can_process(".markdown")
        assert not processor.can_process(".txt")

    def test_extract_text_from_markdown(self, temp_dir):
        """Test text extraction from Markdown file."""
        file_path = os.path.join(temp_dir, "test.md")
        with open(file_path, "w", encoding="utf-8") as f:
            f.write(
                "# Header\n\n"
                "Contact: **John Doe** at john@example.com\n"
                "- Item 1\n- Item 2\n"
            )
        processor = MarkdownProcessor()
        text = processor.extract_text(file_path)
        assert "John Doe" in text
        assert "john@example.com" in text
        assert "Header" in text


class TestSqliteProcessor:
    """Tests for SQLite processor."""

    def test_can_process_sqlite(self):
        """Test that SQLite processor recognizes .db extension."""
        processor = SqliteProcessor()
        assert processor.can_process(".db")
        assert processor.can_process(".sqlite")
        assert processor.can_process(".sqlite3")
        assert not processor.can_process(".txt")

    def test_extract_text_from_sqlite(self, temp_dir):
        """Test text extraction from SQLite database."""
        db_path = os.path.join(temp_dir, "test.db")
        conn = sqlite3.connect(db_path)
        conn.execute("CREATE TABLE contacts (name TEXT, email TEXT)")
        conn.execute("INSERT INTO contacts VALUES ('John Doe', 'john@example.com')")
        conn.commit()
        conn.close()

        processor = SqliteProcessor()
        chunks = list(processor.extract_text(db_path))
        text = " ".join(chunks)
        assert "John Doe" in text
        assert "john@example.com" in text
        assert "contacts" in text.lower()

    def test_extract_text_handles_binary_blob_without_raising(self, temp_dir):
        """Binary BLOB payloads don't crash extraction of the sibling text column.

        latin-1 can decode any byte value, so the "skip undecodable BLOB" branch
        is a defensive fallback rather than something reachable with real bytes;
        this asserts the row (and its non-BLOB column) still comes through.
        """
        db_path = os.path.join(temp_dir, "blob.db")
        conn = sqlite3.connect(db_path)
        conn.execute("CREATE TABLE files (name TEXT, payload BLOB)")
        conn.execute(
            "INSERT INTO files VALUES (?, ?)",
            ("report.pdf", b"\x89PNG\r\n\x1a\n\x00\x01\x02\x03"),
        )
        conn.commit()
        conn.close()

        processor = SqliteProcessor()
        chunks = list(processor.extract_text(db_path))
        text = " ".join(chunks)
        assert "report.pdf" in text


class TestZipProcessor:
    """Tests for ZIP processor."""

    def test_can_process_zip(self):
        """Test that ZIP processor recognizes .zip extension."""
        processor = ZipProcessor()
        assert processor.can_process(".zip")
        assert processor.can_process(".ZIP")
        assert not processor.can_process(".txt")

    def test_extract_text_from_zip(self, temp_dir):
        """Test text extraction from ZIP archive."""
        zip_path = os.path.join(temp_dir, "test.zip")
        with zipfile.ZipFile(zip_path, "w") as zf:
            zf.writestr("doc.txt", "Contact: user@example.com")

        processor = ZipProcessor()
        chunks = list(processor.extract_text(zip_path))
        text = " ".join(chunks)
        assert "user@example.com" in text
        assert "doc.txt" in text

    def test_extract_text_from_zip_skips_directories(self, temp_dir):
        """Test that ZIP processor skips directory entries."""
        zip_path = os.path.join(temp_dir, "test.zip")
        with zipfile.ZipFile(zip_path, "w") as zf:
            zf.writestr("subdir/", "")
            zf.writestr("subdir/file.txt", "content")

        processor = ZipProcessor()
        chunks = list(processor.extract_text(zip_path))
        text = " ".join(chunks)
        assert "content" in text

    def test_extract_text_from_zip_multiple_files(self, temp_dir):
        """Test ZIP with multiple text files."""
        zip_path = os.path.join(temp_dir, "multi.zip")
        with zipfile.ZipFile(zip_path, "w") as zf:
            zf.writestr("a.txt", "File A: alpha@example.com")
            zf.writestr("b.txt", "File B: beta@example.com")

        processor = ZipProcessor()
        chunks = list(processor.extract_text(zip_path))
        text = " ".join(chunks)
        assert "alpha@example.com" in text
        assert "beta@example.com" in text
        assert len(chunks) == 2


class TestXmlProcessor:
    """Tests for XML processor."""

    def test_extract_text_from_xml(self, temp_dir):
        """Test text extraction from XML file."""
        file_path = os.path.join(temp_dir, "test.xml")
        with open(file_path, "w", encoding="utf-8") as f:
            f.write(
                '<?xml version="1.0"?>'
                "<root><name>John Doe</name><email>john@example.com</email></root>"
            )
        processor = XmlProcessor()
        text = processor.extract_text(file_path)
        assert "John Doe" in text
        assert "john@example.com" in text

    def test_extract_text_from_malformed_xml_raises(self, temp_dir):
        """Test that malformed XML raises ParseError (no unsafe regex fallback)."""
        file_path = os.path.join(temp_dir, "malformed.xml")
        with open(file_path, "w", encoding="utf-8") as f:
            f.write(
                "<root><name>John Doe</name><email>john@example.com</email></root"
            )  # missing >
        processor = XmlProcessor()
        # Malformed XML must NOT be silently processed via regex fallback,
        # as that would bypass defusedxml security guarantees.
        with pytest.raises(Exception):
            processor.extract_text(file_path)


class TestVcfProcessor:
    """Tests for VCF processor."""

    def test_can_process_vcf(self):
        """Test that VCF processor recognizes .vcf extension."""
        processor = VcfProcessor()
        assert processor.can_process(".vcf")
        assert processor.can_process(".VCF")
        assert not processor.can_process(".txt")

    def test_extract_text_from_vcf(self, temp_dir):
        """Test text extraction from vCard file."""
        file_path = os.path.join(temp_dir, "test.vcf")
        with open(file_path, "w", encoding="utf-8") as f:
            f.write("BEGIN:VCARD\nFN:John Doe\nEMAIL:john@example.com\nEND:VCARD\n")
        processor = VcfProcessor()
        text = processor.extract_text(file_path)
        assert "John Doe" in text
        assert "john@example.com" in text


class TestIcalProcessor:
    """Tests for iCal processor."""

    def test_can_process_ics(self):
        """Test that iCal processor recognizes .ics extension."""
        processor = IcalProcessor()
        assert processor.can_process(".ics")
        assert processor.can_process(".ICAL")
        assert not processor.can_process(".txt")

    def test_extract_text_from_ics(self, temp_dir):
        """Test text extraction from iCalendar file."""
        file_path = os.path.join(temp_dir, "test.ics")
        with open(file_path, "w", encoding="utf-8") as f:
            f.write(
                "BEGIN:VCALENDAR\n"
                "BEGIN:VEVENT\n"
                "SUMMARY:Meeting with John\n"
                "DESCRIPTION:Contact john@example.com\n"
                "LOCATION:Room 123\n"
                "END:VEVENT\n"
                "END:VCALENDAR\n"
            )
        processor = IcalProcessor()
        text = processor.extract_text(file_path)
        assert "Meeting with John" in text
        assert "john@example.com" in text
        assert "Room 123" in text


class TestPropertiesProcessor:
    """Tests for Properties processor."""

    def test_can_process_properties(self):
        """Test that Properties processor recognizes .properties extension."""
        processor = PropertiesProcessor()
        assert processor.can_process(".properties")
        assert processor.can_process(".ini")
        assert processor.can_process(".cfg")
        assert not processor.can_process(".txt")

    def test_extract_text_from_properties(self, temp_dir):
        """Test text extraction from Java properties file."""
        file_path = os.path.join(temp_dir, "test.properties")
        with open(file_path, "w", encoding="utf-8") as f:
            f.write("user.name=John Doe\napi.key=secret123\n")
        processor = PropertiesProcessor()
        text = processor.extract_text(file_path)
        assert "John Doe" in text
        assert "secret123" in text


class TestMboxProcessor:
    """Tests for MBOX processor."""

    def test_can_process_mbox(self):
        """Test that MBOX processor recognizes .mbox extension."""
        processor = MboxProcessor()
        assert processor.can_process(".mbox")
        assert processor.can_process(".MBOX")
        assert not processor.can_process(".eml")

    def test_extract_text_from_mbox(self, temp_dir):
        """Test text extraction from MBOX mailbox file."""
        file_path = os.path.join(temp_dir, "test.mbox")
        with open(file_path, "wb") as f:
            f.write(
                b"From sender@example.com Mon Jan 01 00:00:00 2024\n"
                b"From: sender@example.com\n"
                b"To: recipient@example.com\n"
                b"Subject: Test\n"
                b"\n"
                b"Body with contact@example.com\n"
            )
        processor = MboxProcessor()
        chunks = list(processor.extract_text(file_path))
        text = " ".join(chunks)
        assert "sender@example.com" in text
        assert "recipient@example.com" in text
        assert "contact@example.com" in text

    def test_can_process_logs_instead_of_silently_swallowing_read_errors(
        self, temp_dir, caplog
    ):
        """can_process's header probe must log, not bare-``pass``, on read failure.

        A directory path reliably raises ``IsADirectoryError`` (an ``OSError``
        subclass) from ``open()``, regardless of process privileges.
        """
        dir_path = os.path.join(temp_dir, "not_a_file")
        os.makedirs(dir_path)

        processor = MboxProcessor()
        with caplog.at_level("DEBUG"):
            result = processor.can_process("", dir_path, "")

        assert result is False
        assert any(
            "Could not read file header" in record.getMessage()
            for record in caplog.records
        )

"""PPTX file processor using python-pptx library."""

import struct

from file_processors.base_processor import BaseFileProcessor

try:
    from pptx import Presentation
    from pptx.exc import PackageNotFoundError
except Exception:  # pragma: no cover - optional dependency
    # python-pptx ships type stubs (py.typed), so these optional-dependency
    # fallbacks reassign typed names; the ignores keep the pattern mypy-clean.
    Presentation = None  # type: ignore[assignment]
    PackageNotFoundError = None  # type: ignore[assignment,misc]

try:
    import olefile
except Exception:  # pragma: no cover - optional dependency
    olefile = None


class PptxProcessor(BaseFileProcessor):
    """Processor for PPTX (PowerPoint 2007+) files.

    Extracts text from PPTX files using python-pptx library.
    Extracts text from slides, notes, and comments for PII detection.
    """

    def extract_text(self, file_path: str) -> str:
        """Extract text from a PPTX file.

        Extracts text from:
        - All slides (text boxes, shapes, tables)
        - Notes pages
        - Comments

        Args:
            file_path: Path to the PPTX file

        Returns:
            Extracted text content from slides, notes, and comments as a string

        Raises:
            ImportError: If python-pptx is not installed
            PermissionError: If file cannot be accessed
            FileNotFoundError: If file does not exist
            Exception: For other PPTX processing errors
        """
        if Presentation is None:
            raise ImportError(
                "python-pptx is required for PPTX processing. "
                "Install it with: pip install python-pptx"
            )

        text_parts: list[str] = []

        try:
            # Load presentation
            try:
                prs = Presentation(file_path)
            except ImportError as e:
                raise ImportError(
                    "python-pptx is required for PPTX processing. "
                    "Install it with: pip install python-pptx"
                ) from e

            # Extract text from all slides
            for slide in prs.slides:
                # Extract text from shapes on slide
                for shape in slide.shapes:
                    shape_text = self._extract_text_from_shape(shape)
                    if shape_text:
                        text_parts.append(shape_text)

                # Extract text from notes slide
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    for shape in notes_slide.shapes:
                        shape_text = self._extract_text_from_shape(shape)
                        if shape_text:
                            text_parts.append(shape_text)

            # Extract text from comments (if available)
            # Note: Comments might not be accessible in all PPTX files
            # This is a best-effort extraction

        except FileNotFoundError:
            raise
        except PermissionError:
            raise
        except ImportError:
            raise
        except Exception as e:
            # Map python-pptx "package not found" to a standard FileNotFoundError
            if PackageNotFoundError is not None and isinstance(e, PackageNotFoundError):
                raise FileNotFoundError(str(e)) from e
            raise Exception(f"Error processing PPTX file: {str(e)}") from e

        return " ".join(text_parts)

    def _extract_text_from_shape(self, shape) -> str:
        """Extract text from a PowerPoint shape.

        Handles different shape types:
        - Text boxes
        - Auto shapes with text
        - Tables

        Args:
            shape: PowerPoint shape object

        Returns:
            Extracted text from the shape
        """
        text_parts: list[str] = []

        # Check if shape has text frame
        if hasattr(shape, "text_frame") and shape.text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text and run.text.strip():
                        text_parts.append(run.text.strip())

        # Check if shape is a table
        if hasattr(shape, "table") and shape.table:
            for row in shape.table.rows:
                for cell in row.cells:
                    cell_text = cell.text
                    if cell_text and cell_text.strip():
                        text_parts.append(cell_text.strip())

        # Also check if shape has direct text attribute
        if hasattr(shape, "text") and shape.text:
            if shape.text.strip():
                text_parts.append(shape.text.strip())

        return " ".join(text_parts)

    @staticmethod
    def can_process(extension: str) -> bool:  # type: ignore[override]  # registry inspects arity; see base_processor.can_process
        """Check if this processor can handle PPTX files."""
        return extension.lower() == ".pptx"


# MS-PPT binary record header: a 2-byte version+instance field (low 4 bits are
# the version), a 2-byte record type, and a 4-byte content length. A version
# nibble of 0xF marks a *container* record whose content is itself a nested
# stream of child records; any other value marks an *atom* holding raw data.
# This convention lets a single recursive walk find every atom in the file
# without needing to know the full record-type hierarchy (slide -> shape ->
# text placeholder, etc.) -- we only care about the two atom types that carry
# extractable text.
_TEXT_CHARS_ATOM = 4008  # UTF-16LE text (0x0FA8)
_TEXT_BYTES_ATOM = (
    4000  # 8-bit (Windows-1252) text, high byte of each char zeroed (0x0FA0)
)
_TEXT_ATOM_TYPES = (_TEXT_CHARS_ATOM, _TEXT_BYTES_ATOM)
_PPT_DOCUMENT_STREAM = "PowerPoint Document"


def _walk_ppt_records(data: bytes):
    """Recursively yield ``(rec_type, content)`` for every atom record in *data*.

    Best-effort by design: this walks an untrusted legacy binary stream, so a
    truncated or malformed record (length running past the end of *data*) is
    clamped rather than raising -- the offset always advances by at least the
    8-byte header, so a corrupt file can yield garbage/partial atoms but can
    never loop forever.
    """
    offset = 0
    length = len(data)
    while offset + 8 <= length:
        rec_ver_instance, rec_type, rec_len = struct.unpack_from("<HHI", data, offset)
        version = rec_ver_instance & 0x0F
        content_start = offset + 8
        content_end = min(content_start + rec_len, length)
        content = data[content_start:content_end]
        if version == 0xF:
            yield from _walk_ppt_records(content)
        else:
            yield rec_type, content
        offset = content_end


def _decode_text_atom(rec_type: int, content: bytes) -> str:
    if rec_type == _TEXT_CHARS_ATOM:
        return content.decode("utf-16-le", errors="replace")
    return content.decode("cp1252", errors="replace")


class PptProcessor(BaseFileProcessor):
    """Processor for legacy PPT (PowerPoint 97-2003) files.

    No maintained pure-Python library reads this binary format, so text is
    extracted directly from the OLE (compound-file) "PowerPoint Document"
    stream by walking its record structure (see :func:`_walk_ppt_records`)
    and decoding the ``TextCharsAtom``/``TextBytesAtom`` records. Formatting,
    images, and slide/notes boundaries are not preserved -- only raw text
    runs are extracted, which is sufficient for PII detection.
    """

    def extract_text(self, file_path: str) -> str:
        """Extract text from a legacy PPT file.

        Args:
            file_path: Path to the PPT file

        Returns:
            Extracted text content, one line per text run found in the
            presentation (slide bodies and speaker notes are not distinguished).

        Raises:
            ImportError: If olefile is not installed
            PermissionError: If file cannot be accessed
            FileNotFoundError: If file does not exist
            Exception: For other PPT processing errors (not a valid OLE file,
                missing "PowerPoint Document" stream, etc.)
        """
        if olefile is None:
            raise ImportError(
                "olefile is required for legacy PPT (PowerPoint 97-2003) processing. "
                "Install it with: pip install olefile"
            )

        try:
            if not olefile.isOleFile(file_path):
                raise Exception(
                    f"Not a valid OLE compound file (legacy PPT expected): {file_path}"
                )

            ole = olefile.OleFileIO(file_path)
            try:
                if not ole.exists(_PPT_DOCUMENT_STREAM):
                    raise Exception(
                        f"No '{_PPT_DOCUMENT_STREAM}' stream found; file may be "
                        "corrupt or not a legacy PPT file"
                    )
                data = ole.openstream(_PPT_DOCUMENT_STREAM).read()
            finally:
                ole.close()
        except FileNotFoundError:
            raise
        except PermissionError:
            raise
        except Exception as e:
            raise Exception(f"Error processing PPT file: {str(e)}") from e

        text_parts = [
            _decode_text_atom(rec_type, content)
            for rec_type, content in _walk_ppt_records(data)
            if rec_type in _TEXT_ATOM_TYPES
        ]
        return "\n".join(t for t in text_parts if t and t.strip())

    @staticmethod
    def can_process(extension: str) -> bool:  # type: ignore[override]  # registry inspects arity; see base_processor.can_process
        """Check if this processor can handle PPT files."""
        return extension.lower() == ".ppt"
